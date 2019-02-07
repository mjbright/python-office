#!/usr/bin/env python3

import sys

files=[
    '1.hello-world.pptx',
    '2.bullets.pptx',
    '3.textbox.pptx',
    '4.addpicture.pptx',
    '5.add_shape.pptx',
    '6.add_table.pptx'
]

from pptx import Presentation

def extract_text(path_pres):
    prs = Presentation(path_pres)

    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []

    slide_no=0
    for slide in prs.slides:
        slide_no+=1
        slide_no_shown=False

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not slide_no_shown:
                        slide_no_shown=True
                        text="slide[{}]: {}".format(slide_no, run.text)
                    else:
                        text=run.text

                    text_runs.append(text)

    print()
    print("FILE: " + path_pres)
    print(text_runs)

if len(sys.argv) > 1:
    files = sys.argv[1:]

for file in files:
    extract_text(file)


